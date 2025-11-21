"""FastAPI main app: Day-1 intake + Day-2 AICMO operator endpoints."""

from __future__ import annotations

import io
import json
import zipfile
from datetime import date, timedelta
from enum import Enum
from typing import Optional

from fastapi import FastAPI, UploadFile, File, HTTPException, Form
from fastapi.responses import PlainTextResponse, JSONResponse, StreamingResponse
from pydantic import BaseModel

from aicmo.io.client_reports import (
    ClientInputBrief,
    AICMOOutputReport,
    MarketingPlanView,
    StrategyPillar,
    CampaignBlueprintView,
    CampaignObjectiveView,
    AudiencePersonaView,
    SocialCalendarView,
    CalendarPostView,
    PerformanceReviewView,
    PerfSummaryView,
    CreativesBlock,
    CreativeRationale,
    ChannelVariant,
    ToneVariant,
    generate_output_report_markdown,
)
from backend.schemas import (
    ClientIntakeForm,
    MarketingPlanReport,
    CampaignBlueprintReport,
    SocialCalendarReport,
    PerformanceReviewReport,
)
from backend.templates import (
    generate_client_intake_text_template,
    generate_blank_marketing_plan_template,
    generate_blank_campaign_blueprint_template,
    generate_blank_social_calendar_template,
    generate_blank_performance_review_template,
)
from backend.pdf_utils import text_to_pdf_bytes
from backend.routers.health import router as health_router


class TemplateFormat(str, Enum):
    json = "json"
    text = "text"
    pdf = "pdf"


class ReportType(str, Enum):
    marketing_plan = "marketing_plan"
    campaign_blueprint = "campaign_blueprint"
    social_calendar = "social_calendar"
    performance_review = "performance_review"


class GenerateRequest(BaseModel):
    brief: ClientInputBrief
    generate_marketing_plan: bool = True
    generate_campaign_blueprint: bool = True
    generate_social_calendar: bool = True
    generate_performance_review: bool = False
    generate_creatives: bool = True


app = FastAPI(title="AICMO API")
app.include_router(health_router, tags=["health"])


# =====================
# INPUT – CLIENT FORM
# =====================


@app.get(
    "/templates/intake",
    response_class=PlainTextResponse,
    summary="Get blank client intake form (text/PDF/JSON schema).",
)
def get_blank_intake_template(fmt: TemplateFormat = TemplateFormat.text):
    if fmt == TemplateFormat.text:
        return generate_client_intake_text_template()
    elif fmt == TemplateFormat.json:
        return JSONResponse(ClientIntakeForm.model_json_schema())
    elif fmt == TemplateFormat.pdf:
        text = generate_client_intake_text_template()
        pdf_bytes = text_to_pdf_bytes(text)
        return StreamingResponse(
            content=iter([pdf_bytes]),
            media_type="application/pdf",
            headers={"Content-Disposition": 'attachment; filename="aicmo_intake_form.pdf"'},
        )
    else:
        raise HTTPException(status_code=400, detail="Unsupported format")


@app.post(
    "/intake/json",
    summary="Submit a filled client intake form as JSON.",
)
def submit_intake_json(payload: ClientIntakeForm):
    return {
        "status": "ok",
        "message": "Intake received",
        "brand_name": payload.brand.brand_name,
    }


@app.post(
    "/intake/file",
    summary="Submit a filled client intake form as text or PDF file.",
)
async def submit_intake_file(file: UploadFile = File(...)):
    content = await file.read()

    if file.content_type == "text/plain":
        return {
            "status": "ok",
            "message": "Text intake received",
            "filename": file.filename,
            "size_bytes": len(content),
        }

    if file.content_type == "application/pdf":
        return {
            "status": "ok",
            "message": "PDF intake received",
            "filename": file.filename,
            "size_bytes": len(content),
        }

    raise HTTPException(
        status_code=400, detail="Only text/plain or application/pdf supported for now"
    )


# =====================
# OUTPUT – BLANK REPORT TEMPLATES
# =====================


@app.get(
    "/templates/report/{report_type}",
    summary="Get blank standardized report template for client-side review.",
)
def get_blank_report_template(
    report_type: ReportType,
    fmt: TemplateFormat = TemplateFormat.text,
):
    if report_type == ReportType.marketing_plan:
        text = generate_blank_marketing_plan_template()
        schema = MarketingPlanReport.model_json_schema()
        filename = "aicmo_marketing_plan"
    elif report_type == ReportType.campaign_blueprint:
        text = generate_blank_campaign_blueprint_template()
        schema = CampaignBlueprintReport.model_json_schema()
        filename = "aicmo_campaign_blueprint"
    elif report_type == ReportType.social_calendar:
        text = generate_blank_social_calendar_template()
        schema = SocialCalendarReport.model_json_schema()
        filename = "aicmo_social_calendar"
    elif report_type == ReportType.performance_review:
        text = generate_blank_performance_review_template()
        schema = PerformanceReviewReport.model_json_schema()
        filename = "aicmo_performance_review"
    else:
        raise HTTPException(status_code=400, detail="Unsupported report type")

    if fmt == TemplateFormat.text:
        return PlainTextResponse(text)
    elif fmt == TemplateFormat.json:
        return JSONResponse(schema)
    elif fmt == TemplateFormat.pdf:
        pdf_bytes = text_to_pdf_bytes(text)
        return StreamingResponse(
            content=iter([pdf_bytes]),
            media_type="application/pdf",
            headers={"Content-Disposition": f'attachment; filename="{filename}.pdf"'},
        )
    else:
        raise HTTPException(status_code=400, detail="Unsupported format")


# =====================
# OUTPUT – ACCEPT FILLED REPORTS (OPTIONAL)
# =====================


@app.post("/reports/marketing_plan")
def submit_marketing_plan(report: MarketingPlanReport):
    return {"status": "ok", "type": "marketing_plan", "brand_name": report.brand_name}


@app.post("/reports/campaign_blueprint")
def submit_campaign_blueprint(report: CampaignBlueprintReport):
    return {
        "status": "ok",
        "type": "campaign_blueprint",
        "campaign_name": report.campaign_name,
    }


@app.post("/reports/social_calendar")
def submit_social_calendar(report: SocialCalendarReport):
    return {"status": "ok", "type": "social_calendar", "brand_name": report.brand_name}


@app.post("/reports/performance_review")
def submit_performance_review(report: PerformanceReviewReport):
    return {"status": "ok", "type": "performance_review", "brand_name": report.brand_name}


# ============================================================
# DAY 2 – AICMO operator endpoints used by the Streamlit UI
# ============================================================


@app.post("/aicmo/generate", response_model=AICMOOutputReport)
def aicmo_generate(req: GenerateRequest):
    """
    Stub generator:
    - Uses brief to build deterministic, client-ready structures.
    - Now includes multi-channel variants, tone variants, email subjects,
      and a creative rationale block.
    """
    b = req.brief.brand
    g = req.brief.goal
    a = req.brief.audience
    s = req.brief.strategy_extras
    today = date.today()

    # MARKETING PLAN
    mp = MarketingPlanView(
        executive_summary=(
            f"{b.brand_name} is aiming to drive {g.primary_goal or 'growth'} "
            f"over the next {g.timeline or 'period'}. This plan covers strategy, "
            "campaign focus, and channel mix."
        ),
        situation_analysis=(
            f"Primary audience: {a.primary_customer}.\n\n"
            "Market context and competition will be refined in future iterations, "
            "but the focus is on consistent, value-driven messaging that compounds over time."
        ),
        strategy=(
            "Position the brand as the default choice for its niche by combining:\n"
            "- consistent social presence\n"
            "- proof-driven storytelling (testimonials, case studies)\n"
            "- clear, repeated core promises across all touchpoints."
        ),
        pillars=[
            StrategyPillar(
                name="Awareness & Reach",
                description="Grow top-of-funnel awareness via social, search and collaborations.",
                kpi_impact="Impressions, reach, profile visits.",
            ),
            StrategyPillar(
                name="Trust & Proof",
                description="Leverage testimonials, case studies and UGC.",
                kpi_impact="Saves, shares, reply DMs, conversion intent.",
            ),
            StrategyPillar(
                name="Conversion & Retention",
                description="Use clear offers, scarcity, and nurture flows to convert and retain.",
                kpi_impact="Leads, trials, purchases, repeat usage.",
            ),
        ],
    )

    # CAMPAIGN BLUEPRINT
    big_idea_industry = b.industry or "your category"
    cb = CampaignBlueprintView(
        big_idea=f"Whenever your ideal buyer thinks of {big_idea_industry}, they remember {b.brand_name} first.",
        objective=CampaignObjectiveView(
            primary=g.primary_goal or "brand_awareness",
            secondary=g.secondary_goal,
        ),
        audience_persona=AudiencePersonaView(
            name="Core Buyer",
            description=(
                f"{a.primary_customer} who is actively looking for better options and wants "
                "less friction, more clarity, and trustworthy proof before committing."
            ),
        ),
    )

    # SOCIAL CALENDAR – simple 7-day stub
    posts: list[CalendarPostView] = []
    for i in range(7):
        d = today + timedelta(days=i)
        theme = "Brand Story" if i == 0 else ("Social Proof" if i == 2 else "Educational")
        posts.append(
            CalendarPostView(
                date=d,
                platform="Instagram",
                theme=theme,
                hook=f"Hook idea for day {i+1}",
                cta="Learn more",
                asset_type="reel" if i % 2 == 0 else "static_post",
                status="planned",
            )
        )

    cal = SocialCalendarView(
        start_date=today,
        end_date=today + timedelta(days=6),
        posts=posts,
    )

    # PERFORMANCE REVIEW – stub if requested
    pr: Optional[PerformanceReviewView] = None
    if req.generate_performance_review:
        pr = PerformanceReviewView(
            summary=PerfSummaryView(
                growth_summary="Performance review will be populated once data is available.",
                wins="- Early engagement signals strong message–market resonance.\n",
                failures="- Limited coverage on secondary channels.\n",
                opportunities="- Double down on top performing content themes and formats.\n",
            )
        )

    # CREATIVES – multi-channel + tones + email subjects + scripts
    creatives: Optional[CreativesBlock] = None
    if req.generate_creatives:
        core_promise = (
            s.success_30_days
            or f"See visible progress towards {g.primary_goal or 'your goal'} within 30 days."
        )

        rationale = CreativeRationale(
            strategy_summary=(
                "The creative system is built around repeating a few clear promises in multiple "
                "formats. Instagram focuses on visual storytelling, LinkedIn focuses on authority "
                "and proof, while X focuses on sharp, scroll-stopping hooks.\n\n"
                "By reusing the same core ideas across platforms, the brand compounds recognition "
                "instead of starting from scratch each time."
            ),
            psychological_triggers=[
                "Social proof",
                "Loss aversion (fear of missing out on better results)",
                "Clarity and specificity (concrete promises)",
                "Authority and expertise",
            ],
            audience_fit=(
                "Ideal for busy decision-makers who scan feeds quickly but respond strongly to "
                "clear proof and repeated, simple promises."
            ),
            risk_notes=(
                "Avoid over-claiming or using fear-heavy framing; keep promises ambitious but "
                "credible and backed by examples whenever possible."
            ),
        )

        # Multi-channel variants
        channel_variants = [
            ChannelVariant(
                platform="Instagram",
                format="reel",
                hook=f"Stop guessing your {big_idea_industry} marketing.",
                caption=(
                    f"Most {big_idea_industry} brands post randomly and hope it works.\n\n"
                    f"{b.brand_name} is switching to a simple, repeatable system: "
                    f"{core_promise}\n\n"
                    "Save this if you're done improvising your growth."
                ),
            ),
            ChannelVariant(
                platform="LinkedIn",
                format="post",
                hook=f"What happened when {b.brand_name} stopped 'posting and praying'.",
                caption=(
                    "We replaced random content with a clear playbook: 3 pillars, 2 offers, "
                    "and 1 simple narrative that repeats everywhere.\n\n"
                    "Result: more consistent leads, fewer 'spray and pray' campaigns."
                ),
            ),
            ChannelVariant(
                platform="X",
                format="thread",
                hook="Most brands don't have a marketing problem. They have a focus problem.",
                caption=(
                    "Thread:\n"
                    "1/ They jump from trend to trend.\n"
                    "2/ They never commit to one clear promise.\n"
                    "3/ Their content feels different on every platform.\n\n"
                    "Fix the focus → the metrics follow."
                ),
            ),
        ]

        # Tone variants
        tone_variants = [
            ToneVariant(
                tone_label="Professional",
                example_caption=(
                    f"{b.brand_name} is implementing a structured, data-aware marketing "
                    "system to replace ad-hoc posting and scattered campaigns."
                ),
            ),
            ToneVariant(
                tone_label="Friendly",
                example_caption=(
                    "No more 'post and pray'. We're building a simple, repeatable marketing "
                    "engine that works even on your busiest weeks."
                ),
            ),
            ToneVariant(
                tone_label="Bold",
                example_caption=(
                    "If your marketing still depends on random ideas and last-minute posts, "
                    "you're leaving serious money on the table."
                ),
            ),
        ]

        # Email subject lines
        email_subject_lines = [
            "Your marketing doesn't need more ideas – it needs a system.",
            f"What happens when {b.brand_name} stops posting randomly?",
            "3 campaigns that can carry your growth for the next 90 days.",
        ]

        # Generic hooks, captions, ad scripts
        hooks = [
            "Stop posting randomly. Start compounding your brand.",
            "Your content is working harder than your strategy. Let's fix that.",
        ]

        captions = [
            "Great marketing is not about doing more. It's about repeating the right things "
            "consistently across channels.",
            "You don't need 100 ideas. You need 5 ideas repeated in 100 smart ways.",
        ]

        scripts = [
            (
                "Opening: Show the chaos (random posts, no clear message).\n"
                "Middle: Introduce the system (3 pillars, 2 offers, 1 narrative).\n"
                "Close: Invite them to take the first step (DM / click / reply)."
            )
        ]

        creatives = CreativesBlock(
            notes="Initial creative system with platform variations, tones, email subjects and scripts.",
            hooks=hooks,
            captions=captions,
            scripts=scripts,
            rationale=rationale,
            channel_variants=channel_variants,
            email_subject_lines=email_subject_lines,
            tone_variants=tone_variants,
        )

    out = AICMOOutputReport(
        marketing_plan=mp,
        campaign_blueprint=cb,
        social_calendar=cal,
        performance_review=pr,
        creatives=creatives,
    )

    return out


@app.post("/aicmo/revise", response_model=AICMOOutputReport)
async def aicmo_revise(
    meta: str = Form(...),
    attachment: Optional[UploadFile] = File(None),
):
    """
    Revision stub:
    - Reads brief + current_output + instructions from 'meta' JSON.
    - For now, just appends a note to the executive summary so you see a change.
    """
    try:
        data = json.loads(meta)
    except json.JSONDecodeError:
        raise HTTPException(status_code=400, detail="Invalid meta JSON")

    current = AICMOOutputReport.model_validate(data["current_output"])

    # Make a shallow copy and tweak executive summary
    mp = current.marketing_plan
    revised_mp = MarketingPlanView(
        executive_summary=mp.executive_summary
        + "\n\n_Revision note: updated according to operator feedback._",
        situation_analysis=mp.situation_analysis,
        strategy=mp.strategy,
        pillars=mp.pillars,
    )

    revised = current.model_copy(update={"marketing_plan": revised_mp})

    return revised


@app.post("/aicmo/export/pdf")
def aicmo_export_pdf(payload: dict):
    """
    Convert markdown (from Streamlit) to a PDF.
    Body: { "markdown": "..." }
    """
    markdown = payload.get("markdown") or ""
    if not markdown.strip():
        raise HTTPException(status_code=400, detail="Markdown content is empty")

    pdf_bytes = text_to_pdf_bytes(markdown)
    return StreamingResponse(
        content=iter([pdf_bytes]),
        media_type="application/pdf",
        headers={"Content-Disposition": 'attachment; filename="aicmo_report.pdf"'},
    )


@app.post("/aicmo/export/pptx")
def aicmo_export_pptx(payload: dict):
    """
    Very simple PPTX generator.
    Body: { "brief": {...}, "output": {...} }
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="python-pptx not installed. Run: pip install python-pptx",
        )

    brief = ClientInputBrief.model_validate(payload["brief"])
    output = AICMOOutputReport.model_validate(payload["output"])

    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"AICMO Report – {brief.brand.brand_name}"
    subtitle = slide.placeholders[1]
    subtitle.text = "Generated by AICMO"

    # Executive summary slide
    mp = output.marketing_plan
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Summary"
    body = slide.placeholders[1].text_frame
    for line in mp.executive_summary.splitlines():
        if not body.text:
            body.text = line
        else:
            body.add_paragraph().text = line

    # Big idea slide
    cb = output.campaign_blueprint
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Campaign Big Idea"
    body = slide.placeholders[1].text_frame
    body.text = cb.big_idea

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    return StreamingResponse(
        content=buf,
        media_type=("application/vnd.openxmlformats-officedocument." "presentationml.presentation"),
        headers={"Content-Disposition": 'attachment; filename="aicmo_report.pptx"'},
    )


@app.post("/aicmo/export/zip")
def aicmo_export_zip(payload: dict):
    """
    Bundle final report and a few text assets into a ZIP.
    Body: { "brief": {...}, "output": {...} }
    """
    brief = ClientInputBrief.model_validate(payload["brief"])
    output = AICMOOutputReport.model_validate(payload["output"])

    report_md = generate_output_report_markdown(brief, output)
    pdf_bytes = text_to_pdf_bytes(report_md)

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("01_Strategy/report.md", report_md)
        z.writestr("01_Strategy/report.pdf", pdf_bytes)
        z.writestr("meta/brand_name.txt", brief.brand.brand_name)

    buf.seek(0)
    return StreamingResponse(
        content=buf,
        media_type="application/zip",
        headers={"Content-Disposition": 'attachment; filename="aicmo_package.zip"'},
    )

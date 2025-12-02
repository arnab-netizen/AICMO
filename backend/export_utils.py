"""
Export utilities for PDF, PPTX, and ZIP generation.

Provides:
- Safe wrappers around export functions with error handling
- Structured error results that can be displayed to operators
- Logging of export failures
- Placeholder detection and blocking before export
"""

import logging
import io
import zipfile
from typing import Dict, Union, Optional

from fastapi.responses import StreamingResponse

from aicmo.io.client_reports import (
    ClientInputBrief,
    AICMOOutputReport,
    generate_output_report_markdown,
)
from aicmo.quality.validators import validate_report, has_blocking_issues
from backend.pdf_utils import text_to_pdf_bytes
from backend.placeholder_utils import report_has_placeholders, format_placeholder_warning
from backend.pdf_renderer import render_agency_pdf, WEASYPRINT_AVAILABLE

logger = logging.getLogger("aicmo.export")


# ──────────────────────────────────────────────────────────────────────────────
# AGENCY PDF ROUTER
# ──────────────────────────────────────────────────────────────────────────────

AGENCY_PACKS = {
    "quick_social_basic",  # quick_social_basic.html
    "strategy_campaign_standard",  # campaign_strategy.html
    # Add more pack keys here as templates are created and upgraded:
    # "full_funnel_growth_suite",
    # "launch_gtm_pack",
    # "brand_turnaround_lab",
    # "retention_crm_booster",
    # "performance_audit_revamp",
}


def should_use_agency_pdf(
    pack_key: str,
    wow_enabled: bool,
    wow_package_key: Optional[str],
) -> bool:
    """
    Decide if we should use the agency (WeasyPrint + template) PDF path.

    Conditions:
    - WOW is enabled
    - wow_package_key is present
    - pack_key has a template in AGENCY_PACKS

    Args:
        pack_key: Pack identifier (e.g., "quick_social_basic")
        wow_enabled: Whether WOW mode is enabled
        wow_package_key: WOW package key from payload

    Returns:
        True if all conditions met, False otherwise
    """
    if not wow_enabled:
        return False
    if not wow_package_key:
        return False
    return pack_key in AGENCY_PACKS


# ──────────────────────────────────────────────────────────────────────────────


class ExportError(Exception):
    """Base exception for export failures."""

    pass


def _validate_report_for_export(
    report: AICMOOutputReport,
) -> Union[None, Dict[str, str]]:
    """
    Validate report against agency-grade standards before exporting to client.

    Returns None if valid, error dict if blocking issues found.

    This is the unified validation entry point that combines both:
    - Old placeholder_utils checks (deprecated but kept for compatibility)
    - New aicmo.quality.validators checks (comprehensive standard)

    Design: Report must pass BOTH checks to be export-ready.
    """
    # First check: new comprehensive validation
    validation_issues = validate_report(report)
    if has_blocking_issues(validation_issues):
        errors = [i for i in validation_issues if i.severity == "error"]
        error_details = "\n".join(
            [f"• {i.section}: {i.message}" for i in errors[:5]]
        )  # Show first 5
        logger.warning(f"Export blocked by validation: {len(errors)} error(s) found")
        return {
            "error": True,
            "message": (
                f"Export blocked: Report validation failed with {len(errors)} issue(s):\n{error_details}"
            ),
            "validation_errors": [{"section": i.section, "message": i.message} for i in errors],
            "blocked_by": "agency_grade_validation",
        }

    # Second check: legacy placeholder detection (backward compatibility)
    placeholder_sections = report_has_placeholders(report)
    if placeholder_sections:
        msg = format_placeholder_warning(placeholder_sections)
        logger.warning(f"Export blocked by legacy placeholder check: {msg}")
        return {
            "error": True,
            "message": f"Export blocked: {msg}",
            "blocked_by": "legacy_placeholder_detection",
        }

    return None


def safe_export_pdf(
    markdown: str, check_placeholders: bool = False
) -> Union[StreamingResponse, Dict[str, str]]:
    """
    Safely convert markdown to PDF with error handling.

    Args:
        markdown: Markdown text to convert.
        check_placeholders: If True, scan markdown for placeholder patterns before export.

    Returns:
        StreamingResponse with PDF content on success.
        Dict with error details on failure (to be returned as JSON error to operator).

    Design: Graceful error return allows caller to show friendly message to operator.
    """
    try:
        if not markdown or not markdown.strip():
            logger.warning("PDF export: empty markdown content provided")
            return {
                "error": True,
                "message": "Cannot export: report is empty. Please generate content first.",
                "export_type": "pdf",
            }

        # Check for placeholders if requested
        if check_placeholders:
            placeholder_patterns = [
                "Hook idea for day",
                "hook idea for day",
                "Performance review will be populated",
                "performance review will be populated",
                "will be populated once data is available",
                "TBD",
                "TO BE DETERMINED",
                "[PLACEHOLDER]",
                "[placeholder]",
            ]

            for pattern in placeholder_patterns:
                if pattern in markdown:
                    logger.warning(f"PDF export blocked: placeholder found: {pattern}")
                    return {
                        "error": True,
                        "message": (
                            f"Export blocked: report contains placeholders. "
                            f"Found: '{pattern}'. Please update or remove these sections before exporting."
                        ),
                        "export_type": "pdf",
                    }

        pdf_bytes = text_to_pdf_bytes(markdown)

        if not pdf_bytes:
            logger.warning("PDF export: text_to_pdf_bytes returned empty result")
            return {
                "error": True,
                "message": "PDF generation produced empty output. Please try again.",
                "export_type": "pdf",
            }

        logger.info(f"PDF export successful ({len(pdf_bytes)} bytes)")

        return StreamingResponse(
            content=iter([pdf_bytes]),
            media_type="application/pdf",
            headers={"Content-Disposition": 'attachment; filename="aicmo_report.pdf"'},
        )

    except ValueError as e:
        # text_to_pdf_bytes raised ValueError (bad markdown)
        logger.error(f"PDF export failed (invalid content): {e}", exc_info=True)
        return {
            "error": True,
            "message": f"Export failed: invalid report content. Error: {str(e)[:100]}",
            "export_type": "pdf",
        }

    except Exception as e:
        # Any other unexpected error
        logger.error(f"PDF export failed (unexpected): {e}", exc_info=True)
        return {
            "error": True,
            "message": "PDF export failed. Please try again or contact support.",
            "export_type": "pdf",
        }


def safe_export_agency_pdf(
    pack_key: str,
    report: Dict,
    wow_enabled: bool,
    wow_package_key: Optional[str],
) -> Optional[bytes]:
    """
    Try to export a PDF using the agency (HTML+CSS) renderer.
    Return None if anything fails so caller can fall back to markdown.

    Args:
        pack_key: Pack identifier (e.g., "quick_social_basic")
        report: Report data dictionary with campaign/brand/objective fields
        wow_enabled: Whether WOW mode is enabled
        wow_package_key: WOW package key (e.g., "quick_social_basic")

    Returns:
        PDF bytes on success, None on any failure

    Design:
        - Only attempts agency rendering if conditions are met
        - Returns None (not an error dict) on any failure
        - Caller is responsible for fallback logic
    """
    print("\n" + "=" * 80)
    print("🎨 AGENCY PDF DEBUG: safe_export_agency_pdf()")
    print(f"   pack_key        = {pack_key}")
    print(f"   wow_enabled     = {wow_enabled}")
    print(f"   wow_package_key = {wow_package_key}")
    print("=" * 80)

    # Check if we should even attempt agency rendering
    if not should_use_agency_pdf(pack_key, wow_enabled, wow_package_key):
        print("🎨 AGENCY PDF DEBUG: conditions not met -> skip agency path")
        print("   - should_use_agency_pdf() returned False")
        return None

    # Check if WeasyPrint is available
    if not WEASYPRINT_AVAILABLE:
        print("🎨 AGENCY PDF DEBUG: WeasyPrint not available -> skip agency path")
        return None

    # Build context for template rendering
    try:
        # Ensure report is a dict (defensive)
        report_data = report if isinstance(report, dict) else {}

        context = {
            "report": report_data,
            "wow_package_key": wow_package_key,
        }

        print(
            f"🎨 AGENCY PDF DEBUG: calling render_agency_pdf() with wow_package_key={wow_package_key}"
        )
        pdf_bytes = render_agency_pdf(context)

        if not pdf_bytes:
            print("🎨 AGENCY PDF DEBUG: render_agency_pdf() returned empty bytes")
            return None

        print(f"✅ AGENCY PDF DEBUG: SUCCESS - Generated {len(pdf_bytes)} bytes")
        logger.info(
            f"Agency PDF export successful ({len(pdf_bytes)} bytes, pack={pack_key}, WOW={wow_package_key})"
        )
        return pdf_bytes

    except Exception as e:
        print(f"❌ AGENCY PDF DEBUG: FAILED in render_agency_pdf(): {e}")
        import traceback

        traceback.print_exc()
        logger.warning(f"Agency PDF rendering failed (pack={pack_key}): {e}")
        return None


def safe_export_pptx(
    brief: ClientInputBrief,
    output: AICMOOutputReport,
    check_placeholders: bool = True,
) -> Union[StreamingResponse, Dict[str, str]]:
    """
    Safely convert brief + output to PPTX with error handling.

    Args:
        brief: Client input brief.
        output: AICMO output report.
        check_placeholders: If True, validate report quality before export.

    Returns:
        StreamingResponse with PPTX content on success.
        Dict with error details on failure.

    Design: Validates input, runs agency-grade checks, attempts PPTX generation, returns graceful error if needed.
    """
    try:
        try:
            from pptx import Presentation
        except ImportError:
            logger.error("PPTX export failed: python-pptx not installed", exc_info=True)
            return {
                "error": True,
                "message": "PPTX export is not available. Please contact support.",
                "export_type": "pptx",
            }

        # Validate input
        if not brief or not output:
            logger.warning("PPTX export: missing brief or output")
            return {
                "error": True,
                "message": "Invalid input: brief or report missing.",
                "export_type": "pptx",
            }

        # Run agency-grade validation if requested
        if check_placeholders:
            validation_error = _validate_report_for_export(output)
            if validation_error:
                logger.warning(f"PPTX export blocked: {validation_error.get('blocked_by')}")
                return validation_error

        mp = output.marketing_plan
        cb = output.campaign_blueprint

        if not mp or not cb:
            logger.warning("PPTX export: missing marketing plan or campaign blueprint")
            return {
                "error": True,
                "message": "Export blocked: report incomplete (missing core sections).",
                "export_type": "pptx",
            }

        # Generate PPTX with expanded slides
        prs = Presentation()

        # Title slide
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = f"AICMO Report – {brief.brand.brand_name}"
            subtitle = slide.placeholders[1]
            subtitle.text = "Generated by AICMO"
        except Exception as e:
            logger.error(f"PPTX export failed: error adding title slide: {e}", exc_info=True)
            return {
                "error": True,
                "message": "Failed to create presentation structure. Please try again.",
                "export_type": "pptx",
            }

        # Executive summary slide
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Executive Summary"
            body = slide.placeholders[1].text_frame
            body.word_wrap = True
            lines = (mp.executive_summary or "").splitlines()
            if lines:
                body.text = lines[0] if lines else "(No summary)"
                for line in lines[1:]:
                    body.add_paragraph().text = line
        except Exception as e:
            logger.error(
                f"PPTX export failed: error adding executive summary slide: {e}", exc_info=True
            )
            logger.debug("PPTX export: continuing without executive summary slide")

        # Situation analysis slide
        try:
            if mp.situation_analysis:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = "Situation Analysis"
                body = slide.placeholders[1].text_frame
                body.word_wrap = True
                body.text = (
                    mp.situation_analysis[:300] + "..."
                    if len(mp.situation_analysis) > 300
                    else mp.situation_analysis
                )
        except Exception as e:
            logger.debug(f"PPTX export: error adding situation analysis slide: {e}")

        # Strategy slide
        try:
            if mp.strategy:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = "Strategy"
                body = slide.placeholders[1].text_frame
                body.word_wrap = True
                body.text = mp.strategy[:300] + "..." if len(mp.strategy) > 300 else mp.strategy
        except Exception as e:
            logger.debug(f"PPTX export: error adding strategy slide: {e}")

        # Strategic pillars slide
        try:
            if mp.pillars:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = "Strategic Pillars"
                body = slide.placeholders[1].text_frame
                body.word_wrap = True
                for pillar in mp.pillars[:3]:  # First 3 pillars
                    pillar_text = f"{pillar.name}: {pillar.description}"[:150] + "..."
                    if not body.text:
                        body.text = pillar_text
                    else:
                        body.add_paragraph().text = pillar_text
        except Exception as e:
            logger.debug(f"PPTX export: error adding pillars slide: {e}")

        # Campaign big idea slide
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Campaign Big Idea"
            body = slide.placeholders[1].text_frame
            body.word_wrap = True
            body.text = cb.big_idea or "(No big idea defined)"
        except Exception as e:
            logger.error(f"PPTX export failed: error adding big idea slide: {e}", exc_info=True)
            logger.debug("PPTX export: continuing without big idea slide")

        # Campaign objective slide
        try:
            if cb.objective:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = "Campaign Objective"
                body = slide.placeholders[1].text_frame
                body.word_wrap = True
                body.text = f"Primary: {cb.objective.primary or 'N/A'}"
                if cb.objective.secondary:
                    body.add_paragraph().text = f"Secondary: {cb.objective.secondary}"
        except Exception as e:
            logger.debug(f"PPTX export: error adding objective slide: {e}")

        # Audience persona slide
        try:
            if cb.audience_persona:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = f"Target Audience: {cb.audience_persona.name}"
                body = slide.placeholders[1].text_frame
                body.word_wrap = True
                body.text = cb.audience_persona.description or "(No description)"
        except Exception as e:
            logger.debug(f"PPTX export: error adding audience slide: {e}")

        # Social calendar overview slide
        try:
            if output.social_calendar and output.social_calendar.posts:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = "Social Content Calendar"
                body = slide.placeholders[1].text_frame
                body.word_wrap = True
                cal = output.social_calendar
                body.text = f"Period: {cal.start_date} to {cal.end_date}"
                body.add_paragraph().text = f"Posts Planned: {len(cal.posts)}"
                body.add_paragraph().text = (
                    f"Platforms: {', '.join(set(p.platform for p in cal.posts[:5]))}"
                )
        except Exception as e:
            logger.debug(f"PPTX export: error adding calendar slide: {e}")

        # KPIs / Measurement slide
        try:
            if mp.messaging_pyramid:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = "Key Messages"
                body = slide.placeholders[1].text_frame
                body.word_wrap = True
                mp_text = mp.messaging_pyramid
                body.text = f"Promise: {mp_text.promise}"
                for msg in mp_text.key_messages[:2]:  # First 2 key messages
                    body.add_paragraph().text = f"• {msg[:100]}" if msg else ""
        except Exception as e:
            logger.debug(f"PPTX export: error adding messaging slide: {e}")

        # Action plan slide
        try:
            if output.action_plan:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = "Action Plan"
                body = slide.placeholders[1].text_frame
                body.word_wrap = True
                ap = output.action_plan
                if ap.quick_wins:
                    body.text = "Quick Wins:"
                    for win in ap.quick_wins[:2]:
                        body.add_paragraph().text = f"• {win[:80]}" if win else ""
                if ap.next_10_days:
                    body.add_paragraph().text = "Next 10 Days:"
                    for item in ap.next_10_days[:2]:
                        body.add_paragraph().text = f"• {item[:80]}" if item else ""
        except Exception as e:
            logger.debug(f"PPTX export: error adding action plan slide: {e}")

        # Creatives overview slide
        try:
            if output.creatives:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = "Creative System"
                body = slide.placeholders[1].text_frame
                body.word_wrap = True
                cr = output.creatives
                if cr.hooks:
                    body.text = "Hooks:"
                    for hook in cr.hooks[:2]:
                        body.add_paragraph().text = f"• {hook[:100]}" if hook else ""
        except Exception as e:
            logger.debug(f"PPTX export: error adding creatives slide: {e}")

        # Write to buffer
        try:
            buf = io.BytesIO()
            prs.save(buf)
            buf.seek(0)
            pptx_bytes = buf.getvalue()

            if not pptx_bytes:
                logger.warning("PPTX export: presentation buffer empty")
                return {
                    "error": True,
                    "message": "PPTX generation produced empty output. Please try again.",
                    "export_type": "pptx",
                }

            logger.info(
                f"PPTX export successful ({len(pptx_bytes)} bytes, {len(prs.slides)} slides)"
            )

            return StreamingResponse(
                content=iter([pptx_bytes]),
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                headers={"Content-Disposition": 'attachment; filename="aicmo_report.pptx"'},
            )
        except Exception as e:
            logger.error(f"PPTX export failed: error writing presentation: {e}", exc_info=True)
            return {
                "error": True,
                "message": "Failed to write presentation file. Please try again.",
                "export_type": "pptx",
            }

    except Exception as e:
        logger.error(f"PPTX export failed (unexpected): {e}", exc_info=True)
        return {
            "error": True,
            "message": "PPTX export failed. Please try again or contact support.",
            "export_type": "pptx",
        }


def safe_export_zip(
    brief: ClientInputBrief,
    output: AICMOOutputReport,
    check_placeholders: bool = True,
) -> Union[StreamingResponse, Dict[str, str]]:
    """
    Safely create ZIP archive with report, personas, creatives.

    Args:
        brief: Client input brief.
        output: AICMO output report.
        check_placeholders: If True, check report for placeholders before export.

    Returns:
        StreamingResponse with ZIP content on success.
        Dict with error details on failure.

    Design: Validates input, generates markdown → PDF → ZIP, returns graceful error if needed.
    """
    try:
        # Validate input
        if not brief or not output:
            logger.warning("ZIP export: missing brief or output")
            return {
                "error": True,
                "message": "Invalid input: brief or report missing.",
                "export_type": "zip",
            }

        # Run agency-grade validation if requested
        if check_placeholders:
            validation_error = _validate_report_for_export(output)
            if validation_error:
                logger.warning(f"ZIP export blocked: {validation_error.get('blocked_by')}")
                return validation_error

        # Generate markdown
        try:
            report_md = generate_output_report_markdown(brief, output)
            if not report_md or not report_md.strip():
                logger.warning("ZIP export: markdown generation produced empty result")
                return {
                    "error": True,
                    "message": "Export blocked: report is empty. Please generate content first.",
                    "export_type": "zip",
                }
        except Exception as e:
            logger.error(f"ZIP export failed: error generating markdown: {e}", exc_info=True)
            return {
                "error": True,
                "message": f"Failed to generate report content. Error: {str(e)[:100]}",
                "export_type": "zip",
            }

        # Generate PDF from markdown
        try:
            pdf_bytes = text_to_pdf_bytes(report_md)
            if not pdf_bytes:
                logger.warning("ZIP export: PDF generation produced empty result")
                # Continue anyway – ZIP should still be created with markdown
                pdf_bytes = b""
        except Exception as e:
            logger.warning(
                f"ZIP export: PDF generation failed (continuing with markdown only): {e}"
            )
            # Non-blocking: ZIP continues without PDF
            pdf_bytes = b""

        # Create ZIP
        try:
            buf = io.BytesIO()
            with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
                # Core strategy
                z.writestr("01_Strategy/report.md", report_md)
                if pdf_bytes:
                    z.writestr("01_Strategy/report.pdf", pdf_bytes)
                z.writestr("meta/brand_name.txt", brief.brand.brand_name or "Unknown")

                # Persona cards
                if output.persona_cards:
                    try:
                        lines = []
                        for p in output.persona_cards:
                            lines.append(f"# {p.name}")
                            lines.append(f"Demographics: {p.demographics}")
                            lines.append(f"Psychographics: {p.psychographics}")
                            lines.append(f"Pain points: {', '.join(p.pain_points)}")
                            lines.append(f"Triggers: {', '.join(p.triggers)}")
                            lines.append(f"Objections: {', '.join(p.objections)}")
                            lines.append(f"Content preferences: {', '.join(p.content_preferences)}")
                            lines.append(f"Primary platforms: {', '.join(p.primary_platforms)}")
                            lines.append(f"Tone preference: {p.tone_preference}")
                            lines.append("")
                        z.writestr("01_Strategy/persona_cards.md", "\n".join(lines))
                    except Exception as e:
                        logger.warning(f"ZIP export: error adding persona cards (continuing): {e}")

                # Creatives multi-format packs
                if output.creatives:
                    try:
                        cr = output.creatives
                        if cr.hooks:
                            z.writestr("02_Creatives/hooks.txt", "\n".join(cr.hooks))
                        if cr.captions:
                            z.writestr("02_Creatives/captions.txt", "\n".join(cr.captions))
                        if cr.scripts:
                            z.writestr(
                                "02_Creatives/scripts.txt",
                                "\n\n---\n\n".join(cr.scripts),
                            )
                        if cr.email_subject_lines:
                            z.writestr(
                                "02_Creatives/email_subject_lines.txt",
                                "\n".join(cr.email_subject_lines),
                            )
                        if cr.cta_library:
                            lines = ["Label | CTA | Context"]
                            for cta in cr.cta_library:
                                lines.append(f"{cta.label} | {cta.text} | {cta.usage_context}")
                            z.writestr("02_Creatives/cta_library.txt", "\n".join(lines))
                        if cr.offer_angles:
                            lines = []
                            for angle in cr.offer_angles:
                                lines.append(f"{angle.label}: {angle.description}")
                                lines.append(f"Example: {angle.example_usage}")
                                lines.append("")
                            z.writestr("02_Creatives/offer_angles.txt", "\n".join(lines))
                    except Exception as e:
                        logger.warning(f"ZIP export: error adding creatives (continuing): {e}")

            buf.seek(0)
            zip_bytes = buf.getvalue()

            if not zip_bytes:
                logger.warning("ZIP export: ZIP buffer empty")
                return {
                    "error": True,
                    "message": "ZIP creation produced empty output. Please try again.",
                    "export_type": "zip",
                }

            logger.info(f"ZIP export successful ({len(zip_bytes)} bytes)")

            return StreamingResponse(
                content=iter([zip_bytes]),
                media_type="application/zip",
                headers={"Content-Disposition": 'attachment; filename="aicmo_package.zip"'},
            )

        except Exception as e:
            logger.error(f"ZIP export failed: error creating ZIP archive: {e}", exc_info=True)
            return {
                "error": True,
                "message": "Failed to create ZIP archive. Please try again or contact support.",
                "export_type": "zip",
            }

    except Exception as e:
        logger.error(f"ZIP export failed (unexpected): {e}", exc_info=True)
        return {
            "error": True,
            "message": "ZIP export failed. Please try again or contact support.",
            "export_type": "zip",
        }

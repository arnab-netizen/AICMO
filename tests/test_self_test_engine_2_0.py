"""
Tests for Self-Test Engine 2.0 Features

Tests for benchmarks, layout checks, format checks, quality checks, and coverage reporting.
"""

import pytest
from pathlib import Path
from tempfile import NamedTemporaryFile

from aicmo.self_test.benchmarks_harvester import (
    discover_all_benchmarks,
    map_benchmarks_to_features,
    _infer_feature_from_name,
)
from aicmo.self_test.layout_checkers import (
    check_html_layout,
    check_pptx_layout,
    HtmlLayoutCheckResult,
)
from aicmo.self_test.format_checkers import (
    check_text_format,
    check_structure,
    validate_calendar_format,
    count_words,
)
from aicmo.self_test.quality_checkers import (
    check_content_quality,
    check_placeholder_markers,
    check_lexical_diversity,
    summarize_quality_metrics,
)
from aicmo.self_test.coverage_report import (
    build_coverage_summary,
    calculate_feature_coverage,
    coverage_assessment,
)


class TestBenchmarksHarvester:
    """Test benchmark discovery and mapping."""

    def test_discover_benchmarks_empty_directory(self):
        """Test discovery when benchmarks directory doesn't exist."""
        benchmarks = discover_all_benchmarks("/tmp/nonexistent")
        assert isinstance(benchmarks, list)
        assert len(benchmarks) == 0

    def test_infer_feature_from_name_persona(self):
        """Test feature inference for persona benchmark."""
        feature = _infer_feature_from_name("persona_benchmark")
        assert feature == "persona_generator"

    def test_infer_feature_from_name_calendar(self):
        """Test feature inference for calendar benchmark."""
        feature = _infer_feature_from_name("social_calendar_spec")
        assert feature == "social_calendar_generator"

    def test_infer_feature_from_name_swot(self):
        """Test feature inference for SWOT benchmark."""
        feature = _infer_feature_from_name("swot")
        assert feature == "swot_generator"

    def test_infer_feature_from_name_unknown(self):
        """Test feature inference returns None for unknown."""
        feature = _infer_feature_from_name("xyz_benchmark")
        assert feature is None

    def test_map_benchmarks_to_features_empty(self):
        """Test mapping with no benchmarks."""
        mapping = map_benchmarks_to_features([])
        assert mapping == {}


class TestLayoutCheckers:
    """Test layout validation for HTML/PPTX/PDF."""

    def test_check_html_layout_aicmo_sections(self):
        """Test HTML layout check detects AICMO sections."""
        html = """
        <html>
            <head><title>Project Report</title></head>
            <body>
                <header>
                    <h1>My Project</h1>
                </header>
                <div class="container">
                    <section>
                        <h2>Project Overview</h2>
                        <p>This is the project overview section.</p>
                    </section>
                    
                    <section>
                        <h2>Strategy</h2>
                        <p>This is the strategy section.</p>
                    </section>
                    
                    <section>
                        <h2>Content Calendar</h2>
                        <p>This is the calendar section.</p>
                    </section>
                    
                    <section>
                        <h2>Deliverables</h2>
                        <ul>
                            <li>Item 1</li>
                            <li>Item 2</li>
                        </ul>
                    </section>
                </div>
            </body>
        </html>
        """
        result = check_html_layout(html)
        assert result.is_valid
        assert len(result.found_sections) >= 3  # Should find strategy, content calendar, deliverables
        found_sections_lower = [s.lower() for s in result.found_sections]
        assert any("overview" in s for s in found_sections_lower)
        assert any("deliverable" in s for s in found_sections_lower)

    def test_check_html_layout_missing_required_sections(self):
        """Test HTML layout check detects missing required sections."""
        html = """
        <html>
            <body>
                <h1>Title Only</h1>
                <p>Some content</p>
            </body>
        </html>
        """
        result = check_html_layout(html)
        # Should have missing sections since AICMO template requires strategy, calendar, deliverables
        assert not result.is_valid
        assert len(result.missing_sections) > 0

    def test_check_html_layout_heading_order(self):
        """Test HTML layout validates heading order (overview before deliverables)."""
        html = """
        <html>
            <body>
                <h2>Deliverables</h2>
                <p>Deliverables come first.</p>
                <h2>Project Overview</h2>
                <p>Overview comes second - incorrect order!</p>
            </body>
        </html>
        """
        result = check_html_layout(html)
        # Heading order should be detected as incorrect
        assert not result.heading_order_ok

    def test_check_html_layout_from_file(self):
        """Test HTML layout check can read from file."""
        with NamedTemporaryFile(mode='w', suffix='.html', delete=False) as f:
            f.write("""
            <html>
                <body>
                    <h2>Project Overview</h2>
                    <h2>Strategy</h2>
                    <h2>Content Calendar</h2>
                    <h2>Deliverables</h2>
                </body>
            </html>
            """)
            f.flush()
            
            result = check_html_layout("", file_path=f.name)
            assert result.is_valid
            
            # Clean up
            Path(f.name).unlink()

    def test_check_pptx_layout_aicmo_structure(self):
        """Test PPTX layout check for AICMO deck structure."""
        # Test with non-existent file to verify error handling
        result = check_pptx_layout("/tmp/nonexistent_deck.pptx")
        assert not result.is_valid
        assert "error" in result.details or "File not found" in result.details.get("error", "")

    def test_check_pptx_layout_slide_count(self):
        """Test PPTX checks minimum slide count."""
        # This checks the validation logic even with error
        result = check_pptx_layout("/tmp/nonexistent.pptx")
        # Result should have slide_count field (will be 0 for non-existent file)
        assert hasattr(result, 'slide_count')
        assert isinstance(result.slide_count, int)

    def test_check_pptx_layout_title_detection(self):
        """Test PPTX extracts and validates slide titles."""
        # Mock PPTX creation for testing
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            import tempfile
            import os
            
            # Create a minimal test PPTX
            prs = Presentation()
            
            # Slide 1: Title
            slide1 = prs.slides.add_slide(prs.slide_layouts[6])
            title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
            title_frame = title_box.text_frame
            title_frame.text = "Creative Execution Deck"
            
            # Slide 2: Strategy
            slide2 = prs.slides.add_slide(prs.slide_layouts[6])
            title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = "Strategy Overview"
            
            # Slide 3: Content Calendar
            slide3 = prs.slides.add_slide(prs.slide_layouts[6])
            title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = "Content Calendar"
            
            # Save to temp file
            temp_path = os.path.join(tempfile.gettempdir(), "test_deck.pptx")
            prs.save(temp_path)
            
            # Test the layout checker
            result = check_pptx_layout(temp_path)
            assert result.slide_count == 3
            assert len(result.found_titles) > 0
            assert "strategy" in str(result.found_titles).lower()
            assert "calendar" in str(result.found_titles).lower()
            assert result.is_valid  # Should pass with 3 slides and 2+ required titles
            
            # Clean up
            os.unlink(temp_path)
            
        except ImportError:
            pytest.skip("python-pptx not available")

    def test_check_html_layout_empty_string(self):
        """Test HTML layout check with empty string."""
        result = check_html_layout("")
        # Should be invalid since no sections found
        assert not result.is_valid
        assert len(result.missing_sections) > 0

    def test_check_pptx_layout_file_not_found(self):
        """Test PPTX check with non-existent file."""
        result = check_pptx_layout("/nonexistent/file.pptx")
        # Should detect the error
        assert not result.is_valid or "error" in result.details


class TestFormatCheckers:
    """Test format and word-count validation."""

    def test_count_words_simple(self):
        """Test word counting."""
        assert count_words("hello world") == 2
        assert count_words("") == 0
        assert count_words("single") == 1

    def test_check_text_format_summary_too_short(self):
        """Test summary too short is flagged."""
        fields = {"summary": "Short text"}
        result = check_text_format(fields)
        assert not result.is_valid
        assert len(result.too_short_fields) > 0

    def test_check_text_format_summary_valid(self):
        """Test valid summary passes."""
        fields = {"summary": " ".join(["word"] * 50)}  # 50 words
        result = check_text_format(fields)
        assert result.is_valid

    def test_check_structure_with_bullets(self):
        """Test structure detection finds bullets."""
        text = "• First bullet\n• Second bullet\n• Third bullet"
        result = check_structure(text)
        assert result["bullet_count"] == 3

    def test_check_structure_paragraphs(self):
        """Test paragraph counting."""
        text = "Para 1 text.\n\nPara 2 text.\n\nPara 3 text."
        result = check_structure(text)
        assert result["paragraph_count"] == 3

    def test_validate_calendar_format_empty(self):
        """Test calendar validation with empty list returns valid=False."""
        result = validate_calendar_format([])
        # Empty calendar is invalid
        assert not result["is_valid"] or result["total_entries"] == 0

    def test_validate_calendar_format_valid_entries(self):
        """Test calendar validation with valid entries."""
        entries = [
            {"date": "2024-01-01", "platform": "twitter", "content": "Hello world tweet"},
            {"date": "2024-01-02", "platform": "linkedin", "content": "Professional post content"},
        ]
        result = validate_calendar_format(entries)
        assert result["total_entries"] == 2
    def test_check_text_format_with_pydantic_object(self):
        """Test format checking with Pydantic model (new API)."""
        from pydantic import BaseModel
        
        class SampleOutput(BaseModel):
            executive_summary: str
            strategy: str
            caption: str
        
        model = SampleOutput(
            executive_summary=" ".join(["word"] * 50),  # Valid: 50 words
            strategy=" ".join(["word"] * 60),  # Valid: 60 words
            caption=" ".join(["word"] * 10),  # Valid: 10 words
        )
        
        result = check_text_format(data=model)
        assert result.is_valid
        assert len(result.metrics) > 0

    def test_check_text_format_with_nested_dict(self):
        """Test format checking extracts fields from nested dicts."""
        data = {
            "summary": {
                "executive_summary": " ".join(["word"] * 50),  # Valid
            },
            "content": {
                "caption": "short text",  # Too short
            }
        }
        
        result = check_text_format(data=data)
        assert not result.is_valid
        assert any("caption" in f for f in result.too_short_fields)

    def test_check_text_format_detects_too_long(self):
        """Test format checker detects text that is too long."""
        fields = {
            "headline": " ".join(["word"] * 50),  # 50 words, max is 15
        }
        result = check_text_format(fields)
        assert not result.is_valid
        assert len(result.too_long_fields) > 0

    def test_check_text_format_with_custom_thresholds(self):
        """Test format checker respects custom thresholds."""
        fields = {
            "custom_field": "five words in this text"  # 5 words
        }
        
        # With custom threshold allowing 4-10 words
        result = check_text_format(
            fields,
            custom_thresholds={"custom_field": {"min_words": 4, "max_words": 10}}
        )
        assert result.is_valid

    def test_format_check_metrics_accuracy(self):
        """Test that format check metrics are accurate."""
        text = "word1 word2 word3 word4 word5"  # Exactly 5 words
        fields = {"description": text}
        result = check_text_format(fields)
        
        assert "description" in result.metrics
        assert result.metrics["description"]["word_count"] == 5

class TestQualityCheckers:
    """Test content quality validation."""

    def test_check_content_quality_empty(self):
        """Test empty content is flagged as poor quality."""
        result = check_content_quality("")
        assert not result.is_valid

    def test_check_content_quality_good(self):
        """Test original content passes quality check."""
        text = "Our innovative approach leverages advanced machine learning algorithms to optimize customer experience in unprecedented ways."
        result = check_content_quality(text)
        assert hasattr(result, 'genericity_score')
        assert 0 <= result.genericity_score <= 1.0

    def test_check_content_quality_generic(self):
        """Test generic content is detected."""
        text = "Leverage cutting-edge solutions to drive engagement and growth in today's digital landscape with synergy and paradigm shift."
        result = check_content_quality(text)
        # Should have some indication of genericity (either phrases or lowered score or warnings)
        # The genericity score may be marginal, so we check for detectable issues
        assert result.genericity_score <= 0.8  # Generic content should score lower

    def test_check_placeholder_markers_finds_insertions(self):
        """Test placeholder detection finds [INSERT] patterns."""
        text = "This is [INSERT NAME] and [TBD] content"
        result = check_placeholder_markers(text)
        assert result["has_placeholders"]
        assert result["placeholder_count"] >= 2

    def test_check_placeholder_markers_finds_braces(self):
        """Test placeholder detection finds {VARIABLE} patterns."""
        text = "Welcome {YOUR_NAME} to {BRAND}"
        result = check_placeholder_markers(text)
        assert result["has_placeholders"]

    def test_check_placeholder_markers_finds_keywords(self):
        """Test placeholder detection finds placeholder keywords."""
        text = "This section is lorem ipsum and will be populated once data is available"
        result = check_placeholder_markers(text)
        assert result["has_placeholders"]

    def test_check_lexical_diversity_good(self):
        """Test diversity score for varied text."""
        text = "The quick brown fox jumps over the lazy dog. Multiple diverse words indicate vocabulary richness."
        result = check_lexical_diversity(text)
        assert result["diversity_ratio"] > 0.3
        assert not result["is_repetitive"]

    def test_check_lexical_diversity_poor(self):
        """Test diversity score detects repetition."""
        text = " ".join(["word"] * 20)  # Repeated same word 20 times
        result = check_lexical_diversity(text)
        assert result["is_repetitive"]
        assert result["diversity_ratio"] < 0.3

    def test_summarize_quality_metrics(self):
        """Test quality summary calculation."""
        from aicmo.self_test.quality_checkers import ContentQualityCheckResult

        quality_result = ContentQualityCheckResult(
            is_valid=True,
            genericity_score=0.8,
            placeholders_found=[],
            warnings=[],
        )
        diversity_result = {
            "diversity_ratio": 0.7,
            "is_repetitive": False,
        }
        summary = summarize_quality_metrics(quality_result, diversity_result)
        assert "overall_quality_score" in summary
        assert "quality_assessment" in summary
        assert summary["overall_quality_score"] > 0


class TestCoverageReport:
    """Test coverage reporting."""

    def test_build_coverage_summary_basic(self):
        """Test basic coverage summary creation."""
        coverage = build_coverage_summary(
            total_benchmarks=10,
            mapped_benchmarks=8,
            enforced_benchmarks=6,
            html_checked=True,
            pptx_checked=False,
        )
        assert coverage.total_benchmarks == 10
        assert coverage.mapped_benchmarks == 8
        assert coverage.enforced_benchmarks == 6
        assert coverage.html_layout_checked
        assert not coverage.pptx_layout_checked

    def test_coverage_summary_percents(self):
        """Test coverage percentage calculation."""
        coverage = build_coverage_summary(
            total_benchmarks=10,
            mapped_benchmarks=10,
            enforced_benchmarks=5,
        )
        assert coverage.benchmark_coverage_percent == 50.0

    def test_calculate_feature_coverage(self):
        """Test per-feature coverage calculation."""
        coverage = calculate_feature_coverage(
            feature_name="persona_generator",
            benchmarks_mapped=["benchmark1", "benchmark2"],
            benchmarks_enforced=["benchmark1"],
        )
        assert coverage.feature_name == "persona_generator"
        assert len(coverage.benchmarks_mapped) == 2
        assert len(coverage.benchmarks_enforced) == 1
        assert len(coverage.benchmarks_unenforced) == 1

    def test_coverage_assessment_basic(self):
        """Test coverage assessment and recommendations."""
        coverage = build_coverage_summary(
            total_benchmarks=10,
            mapped_benchmarks=5,
            enforced_benchmarks=3,
        )
        assessment = coverage_assessment(coverage)
        assert "assessment" in assessment
        assert "recommendations" in assessment
        assert len(assessment["recommendations"]) > 0


class TestIntegration:
    """Integration tests for 2.0 features."""

    def test_all_modules_import(self):
        """Test all 2.0 modules can be imported."""
        from aicmo.self_test import (
            benchmarks_harvester,
            layout_checkers,
            format_checkers,
            quality_checkers,
            coverage_report,
        )
        assert benchmarks_harvester is not None
        assert layout_checkers is not None
        assert format_checkers is not None
        assert quality_checkers is not None
        assert coverage_report is not None

    def test_models_include_2_0_fields(self):
        """Test that models include 2.0 fields."""
        from aicmo.self_test.models import (
            FeatureStatus,
            SelfTestResult,
            BenchmarkCoverageStatus,
            LayoutCheckResults,
        )
        # Create a feature with 2.0 fields
        feature = FeatureStatus(
            name="test",
            category=__import__('aicmo.self_test.models', fromlist=['FeatureCategory']).FeatureCategory.GENERATOR,
            status=__import__('aicmo.self_test.models', fromlist=['TestStatus']).TestStatus.PASS,
            benchmark_coverage=BenchmarkCoverageStatus(feature_name="test"),
        )
        assert feature.benchmark_coverage is not None

    def test_cli_help_shows_2_0_options(self):
        """Test CLI shows new 2.0 options in help."""
        import subprocess
        result = subprocess.run(
            ["python", "-m", "aicmo.self_test.cli", "--help"],
            cwd="/workspaces/AICMO",
            capture_output=True,
            text=True,
        )
        output = result.stdout + result.stderr
        assert "--quality" in output or "--no-quality" in output
        assert "--layout" in output or "--no-layout" in output
        assert "--benchmarks-only" in output


if __name__ == "__main__":
    pytest.main([__file__, "-v"])

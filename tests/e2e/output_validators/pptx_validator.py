"""PPTX validator."""

import os
from typing import Dict, Any
from .base import BaseValidator


class PPTXValidator(BaseValidator):
    """Validates PPTX presentations."""
    
    def validate_structure(self, filepath: str, contract: Dict) -> Dict[str, Any]:
        """Validate PPTX structure."""
        result = {
            'valid': True,
            'errors': [],
            'warnings': []
        }
        
        if not os.path.exists(filepath):
            result['valid'] = False
            result['errors'].append(f"PPTX file not found: {filepath}")
            return result
        
        try:
            from pptx import Presentation
            
            try:
                prs = Presentation(filepath)
                
                # Check slide count
                slide_count = len(prs.slides)
                min_slides = contract.get('min_slides', 0)
                max_slides = contract.get('max_slides', 9999)
                
                if slide_count < min_slides:
                    result['valid'] = False
                    result['errors'].append(
                        f"Slide count {slide_count} < minimum {min_slides}"
                    )
                
                if slide_count > max_slides:
                    result['valid'] = False
                    result['errors'].append(
                        f"Slide count {slide_count} > maximum {max_slides}"
                    )
                
                # Check for speaker notes if required
                if contract.get('speaker_notes_required', False):
                    min_notes_length = contract.get('min_speaker_notes_per_slide', 0)
                    slides_without_notes = 0
                    
                    for slide in prs.slides:
                        notes_text = slide.notes_slide.notes_text_frame.text
                        if len(notes_text.strip()) < min_notes_length:
                            slides_without_notes += 1
                    
                    if slides_without_notes > 0:
                        result['warnings'].append(
                            f"{slides_without_notes} slides have insufficient speaker notes"
                        )
                
                result['slide_count'] = slide_count
                result['has_content'] = slide_count > 0
                
            except Exception as e:
                result['valid'] = False
                result['errors'].append(f"PPTX read error: {str(e)}")
                
        except ImportError:
            result['warnings'].append("python-pptx not installed, skipping PPTX validation")
        
        return result
    
    def validate_safety(self, filepath: str, contract: Dict) -> Dict[str, Any]:
        """Validate PPTX safety."""
        result = {
            'valid': True,
            'errors': [],
            'warnings': []
        }
        
        # Extract text
        text = self.extract_text(filepath)
        
        if not text:
            result['warnings'].append("No text extracted from PPTX")
            return result
        
        # Check placeholders
        placeholder_valid, placeholder_issues = self.validate_placeholders(text, contract)
        if not placeholder_valid:
            result['valid'] = False
            result['errors'].extend(placeholder_issues)
        
        # Check required strings
        required_valid, required_issues = self.validate_required_strings(text, contract)
        if not required_valid:
            result['valid'] = False
            result['errors'].extend(required_issues)
        
        # Check file size
        size_valid, size_issues = self.validate_file_size(filepath, contract)
        if not size_valid:
            result['valid'] = False
            result['errors'].extend(size_issues)
        
        result['no_placeholders'] = placeholder_valid
        result['required_strings_present'] = required_valid
        
        return result
    
    def extract_text(self, filepath: str) -> str:
        """Extract text from PPTX."""
        try:
            from pptx import Presentation
            
            text_parts = []
            prs = Presentation(filepath)
            
            for slide in prs.slides:
                # Extract from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
                
                # Extract from notes
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text:
                    text_parts.append(notes_text)
            
            return '\n'.join(text_parts)
        except (ImportError, Exception):
            return ""

"""
PPTX Renderer for AICMO Delivery Packages

Uses python-pptx to generate professional client-facing presentations.
"""

import os
from typing import Dict, Any
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


def render_delivery_pptx(
    path: str,
    manifest: Dict[str, Any],
    artifacts: Dict[str, Any],
    branding: Dict[str, Any]
) -> str:
    """
    Generate PPTX delivery presentation.
    
    Args:
        path: Output file path
        manifest: Delivery manifest dict
        artifacts: Dict of artifact_type -> Artifact
        branding: Branding configuration
    
    Returns:
        Path to generated PPTX
    """
    if not HAS_PPTX:
        # Fallback: create empty file with error message
        with open(path, 'w') as f:
            f.write("ERROR: python-pptx not installed. Run: pip install python-pptx")
        return path
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    _add_title_slide(prs, manifest, branding)
    
    # Slide 2: Agenda
    _add_agenda_slide(prs, artifacts, branding)
    
    # Content slides
    if artifacts.get("intake"):
        _add_intake_slide(prs, artifacts["intake"], branding)
    
    if artifacts.get("strategy"):
        _add_strategy_slides(prs, artifacts["strategy"], branding)
    
    if artifacts.get("creatives"):
        _add_creatives_slide(prs, artifacts["creatives"], branding)
    
    if artifacts.get("execution"):
        _add_execution_slide(prs, artifacts["execution"], branding)
    
    if artifacts.get("monitoring"):
        _add_monitoring_slide(prs, artifacts["monitoring"], branding)
    
    # Save
    prs.save(path)
    
    return path


def _add_title_slide(prs, manifest, branding):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Campaign Delivery Report"
    
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    top = Inches(3.8)
    height = Inches(1.5)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    campaign_id = manifest["ids"].get("campaign_id", "Unknown")
    client_id = manifest["ids"].get("client_id", "Unknown")
    
    tf.text = f"Campaign: {campaign_id}\nClient: {client_id}\n{datetime.utcnow().strftime('%Y-%m-%d')}"
    
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Footer
    top = Inches(6.5)
    height = Inches(0.5)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = f"Prepared by {branding.get('agency_name', 'AICMO')}"
    
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.CENTER


def _add_agenda_slide(prs, artifacts, branding):
    """Add agenda slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
    
    title = slide.shapes.title
    title.text = "Agenda"
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    items = []
    if artifacts.get("intake"):
        items.append("Intake Summary")
    if artifacts.get("strategy"):
        items.append("Strategy Contract (8 Layers)")
    if artifacts.get("creatives"):
        items.append("Creatives Summary")
    if artifacts.get("execution"):
        items.append("Execution Plan")
    if artifacts.get("monitoring"):
        items.append("Monitoring Setup")
    
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(20)


def _add_intake_slide(prs, intake, branding):
    """Add intake slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Intake Summary"
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    content = intake.content
    
    fields = [
        f"Client: {content.get('client_name', 'N/A')}",
        f"Industry: {content.get('industry', 'N/A')}",
        f"Objective: {content.get('objective', 'N/A')}",
        f"Primary Offer: {content.get('primary_offer', 'N/A')}",
        f"Target Audience: {content.get('target_audience', 'N/A')}"
    ]
    
    for field in fields:
        p = tf.add_paragraph()
        p.text = field
        p.font.size = Pt(16)
        p.level = 0


def _add_strategy_slides(prs, strategy, branding):
    """Add strategy slides (one per layer)"""
    content = strategy.content
    
    # Overview slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Strategy Contract"
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "8-Layer Strategic Framework"
    p.font.size = Pt(18)
    p.font.bold = True
    
    layers = [
        "1. Business Reality Alignment",
        "2. Market & Competitive Truth",
        "3. Audience Decision Psychology",
        "4. Value Proposition Architecture",
        "5. Strategic Narrative",
        "6. Channel Strategy",
        "7. Execution Constraints",
        "8. Measurement & Learning Loop"
    ]
    
    for layer in layers:
        p = tf.add_paragraph()
        p.text = layer
        p.font.size = Pt(14)
        p.level = 0
    
    # Layer detail slides (simplified - show key fields only)
    if "layer1_business_reality" in content:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Layer 1: Business Reality"
        
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        
        layer1 = content["layer1_business_reality"]
        p = tf.add_paragraph()
        p.text = f"Business Model: {layer1.get('business_model_summary', 'N/A')}"
        p.font.size = Pt(16)
        
        p = tf.add_paragraph()
        p.text = f"Bottleneck: {layer1.get('bottleneck', 'N/A')}"
        p.font.size = Pt(16)


def _add_creatives_slide(prs, creatives, branding):
    """Add creatives slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Creatives Summary"
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    content = creatives.content
    
    p = tf.add_paragraph()
    p.text = f"Status: {creatives.status.value}"
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = f"Version: {creatives.version}"
    p.font.size = Pt(16)
    
    if "brief" in content and isinstance(content["brief"], dict):
        theme = content["brief"].get("campaign_theme", "N/A")
        p = tf.add_paragraph()
        p.text = f"Campaign Theme: {theme}"
        p.font.size = Pt(16)


def _add_execution_slide(prs, execution, branding):
    """Add execution slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Execution Plan"
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    content = execution.content
    
    p = tf.add_paragraph()
    p.text = f"Status: {execution.status.value}"
    p.font.size = Pt(16)
    
    if "timeline" in content:
        p = tf.add_paragraph()
        p.text = f"Timeline: {content['timeline']}"
        p.font.size = Pt(16)


def _add_monitoring_slide(prs, monitoring, branding):
    """Add monitoring slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "Monitoring Setup"
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    content = monitoring.content
    
    p = tf.add_paragraph()
    p.text = f"Status: {monitoring.status.value}"
    p.font.size = Pt(16)
    
    if "kpi_config" in content and isinstance(content["kpi_config"], dict):
        kpis = content["kpi_config"].get("selected_kpis", [])
        if kpis:
            p = tf.add_paragraph()
            p.text = f"KPIs: {', '.join(kpis)}"
            p.font.size = Pt(16)

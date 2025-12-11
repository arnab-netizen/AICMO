"""
Layout and Structure Checkers for Self-Test Engine 2.0

Validates the structure and layout of generated client-facing outputs:
- HTML summaries (section presence, heading order)
- PPTX decks (slide count, required slide titles)
- PDF exports (page count, title presence) - stub if parser not available

This module provides concrete, evidence-based validation:
- Never claims to check things we're not actually parsing
- Reports actual counts and findings
- Explicit about unimplemented checks
"""

import logging
import os
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)


@dataclass
class HtmlLayoutCheckResult:
    """Result of HTML layout validation."""

    is_valid: bool
    missing_sections: List[str] = field(default_factory=list)
    found_sections: List[str] = field(default_factory=list)
    heading_order_ok: bool = True
    found_headings: List[str] = field(default_factory=list)
    details: Dict[str, Any] = field(default_factory=dict)


@dataclass
class PptxLayoutCheckResult:
    """Result of PPTX layout validation."""

    is_valid: bool
    slide_count: int = 0
    required_titles_present: List[str] = field(default_factory=list)
    missing_titles: List[str] = field(default_factory=list)
    found_titles: List[str] = field(default_factory=list)
    details: Dict[str, Any] = field(default_factory=dict)


@dataclass
class PdfLayoutCheckResult:
    """Result of PDF layout validation."""

    is_valid: bool
    page_count: int = 0
    has_title_on_first_page: bool = False
    details: Dict[str, Any] = field(default_factory=dict)


def check_html_layout(html_content: str = "", file_path: Optional[str] = None) -> HtmlLayoutCheckResult:
    """
    Validate HTML layout and structure.
    
    Checks for expected sections based on output_packager.py template:
    - Project Overview
    - Strategy
    - Content Calendar
    - Deliverables
    
    Also validates heading order (overview before deliverables, etc.)

    Args:
        html_content: HTML string to validate
        file_path: If provided, file is read from this path instead

    Returns:
        HtmlLayoutCheckResult with validation details
    """
    result = HtmlLayoutCheckResult(is_valid=False)

    if file_path:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                html_content = f.read()
        except Exception as e:
            logger.error(f"Could not read HTML file: {e}")
            result.details["error"] = str(e)
            return result

    # Expected sections in AICMO HTML reports (from output_packager.py template)
    expected_sections = {
        "project overview": "Project Overview section",
        "strategy": "Strategy section",
        "content calendar": "Content Calendar section",
        "deliverables": "Deliverables section",
    }

    try:
        from html.parser import HTMLParser

        class HeadingExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self.headings = []
                self.in_heading = False
                self.current_heading = ""

            def handle_starttag(self, tag, attrs):
                if tag in ["h1", "h2", "h3", "h4", "h5", "h6"]:
                    self.in_heading = True
                    self.current_heading = ""

            def handle_endtag(self, tag):
                if tag in ["h1", "h2", "h3", "h4", "h5", "h6"]:
                    self.in_heading = False
                    if self.current_heading.strip():
                        self.headings.append(self.current_heading.strip().lower())

            def handle_data(self, data):
                if self.in_heading:
                    self.current_heading += data

        extractor = HeadingExtractor()
        extractor.feed(html_content)
        found_headings = extractor.headings
        result.found_headings = found_headings

        # Check for expected sections
        missing = []
        found_secs = []
        for expected_key, expected_label in expected_sections.items():
            found = any(expected_key in heading for heading in found_headings)
            if found:
                found_secs.append(expected_key)
            else:
                missing.append(expected_label)

        result.missing_sections = missing
        result.found_sections = found_secs

        # Check heading order (overview should come before deliverables)
        has_overview = any("overview" in h for h in found_headings)
        has_deliverables = any("deliverable" in h for h in found_headings)

        if has_overview and has_deliverables:
            overview_idx = next(i for i, h in enumerate(found_headings) if "overview" in h)
            deliverables_idx = next(i for i, h in enumerate(found_headings) if "deliverable" in h)
            result.heading_order_ok = overview_idx < deliverables_idx
        else:
            result.heading_order_ok = True  # Can't check order if sections missing

        # Determine validity: all expected sections present and good order
        result.is_valid = (len(missing) == 0) and result.heading_order_ok

        # Add details
        result.details["total_headings_found"] = len(found_headings)
        result.details["expected_sections"] = list(expected_sections.keys())
        result.details["has_overview"] = has_overview
        result.details["has_deliverables"] = has_deliverables

        logger.info(
            f"HTML layout check: {len(found_headings)} headings found, "
            f"{len(missing)} missing sections"
        )

    except Exception as e:
        logger.error(f"Error parsing HTML: {e}")
        result.details["error"] = str(e)

    return result


def check_pptx_layout(file_path: str) -> PptxLayoutCheckResult:
    """
    Validate PowerPoint deck layout and structure.
    
    Checks for expected slide titles based on output_packager.py implementation:
    - Strategy-related slide
    - Content/platform slide
    - Calendar-related slide
    
    Also validates minimum slide count (>= 3).

    Args:
        file_path: Path to PPTX file

    Returns:
        PptxLayoutCheckResult with validation details
    """
    result = PptxLayoutCheckResult(is_valid=False)

    # Expected slide title patterns in AICMO decks
    # (based on output_packager.py implementation)
    expected_slide_patterns = [
        "strategy",  # "Strategy Overview"
        "content",   # "[Platform] Content" or generic Content
        "calendar",  # "Content Calendar"
    ]

    if not os.path.exists(file_path):
        logger.error(f"PPTX file not found: {file_path}")
        result.details["error"] = f"File not found: {file_path}"
        return result

    try:
        from pptx import Presentation
    except ImportError:
        logger.error("python-pptx not available")
        result.details["error"] = "python-pptx not available"
        return result

    try:
        prs = Presentation(file_path)
        slide_count = len(prs.slides)
        result.slide_count = slide_count

        # Extract all slide titles
        all_titles = []
        for slide_idx, slide in enumerate(prs.slides):
            slide_title = ""

            # Try to find title in shapes
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text and len(text) < 150:  # Prioritize shorter text (usually titles)
                        slide_title = text
                        break

            if slide_title:
                all_titles.append(slide_title)

        result.found_titles = all_titles

        # Check for expected slide patterns
        missing_titles = []
        found_required = []
        all_titles_lower = [t.lower() for t in all_titles]

        for expected in expected_slide_patterns:
            found = any(expected in title for title in all_titles_lower)
            if found:
                found_required.append(expected)
            else:
                missing_titles.append(expected)

        result.missing_titles = missing_titles
        result.required_titles_present = found_required

        # Validate structure
        # AICMO decks should have:
        # - At least 3 slides (title, strategy, content/calendar)
        # - At least 2 of 3 required titles present
        min_slides_ok = slide_count >= 3
        required_titles_ok = len(found_required) >= 2

        result.is_valid = min_slides_ok and required_titles_ok

        # Add details
        result.details["total_slides"] = slide_count
        result.details["expected_slide_patterns"] = expected_slide_patterns
        result.details["min_slides_required"] = 3
        result.details["min_required_titles"] = 2

        logger.info(
            f"PPTX layout check: {slide_count} slides, "
            f"{len(found_required)}/{len(expected_slide_patterns)} required titles found"
        )

    except Exception as e:
        logger.error(f"Error reading PPTX: {e}")
        result.details["error"] = str(e)

    return result


def check_pdf_layout(file_path: str) -> PdfLayoutCheckResult:
    """
    Validate PDF layout and structure.

    Checks for:
    - Page count >= minimum (1)
    - Title/heading on first page
    
    Note: PDF parsing requires additional libraries (pdfplumber, PyPDF2) that are
    not currently in requirements.txt. Only reportlab (for PDF generation) is available.
    This returns a stub result with is_valid=False and explicit reason if no PDF
    parser is available.

    Args:
        file_path: Path to PDF file

    Returns:
        PdfLayoutCheckResult with validation findings
    """
    result = PdfLayoutCheckResult(is_valid=False)

    if not os.path.exists(file_path):
        logger.warning(f"PDF file not found: {file_path}")
        result.details["error"] = "File not found"
        result.details["reason"] = "PDF layout check not implemented"
        return result

    # Try to import a PDF parser library
    pdf_parser_available = False
    pdf_lib = None

    try:
        import PyPDF2
        pdf_lib = "PyPDF2"
        pdf_parser_available = True
    except ImportError:
        pass

    if not pdf_parser_available:
        try:
            import pdfplumber
            pdf_lib = "pdfplumber"
            pdf_parser_available = True
        except ImportError:
            pass

    if not pdf_parser_available:
        logger.warning(
            "No PDF parsing library available (pdfplumber or PyPDF2 required)"
        )
        result.details["reason"] = "PDF layout check not implemented"
        result.details["required_libraries"] = ["pdfplumber", "PyPDF2"]
        result.details[
            "message"
        ] = "No PDF parsing library available. Add pdfplumber or PyPDF2 to requirements to enable PDF layout checks."
        return result

    # Try to read PDF with available library
    try:
        if pdf_lib == "PyPDF2":
            import PyPDF2

            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                page_count = len(reader.pages)
                result.page_count = page_count

                # Check first page for title
                if page_count > 0:
                    first_page = reader.pages[0]
                    text = first_page.extract_text() or ""
                    # Simple heuristic: first page should have some text
                    result.has_title_on_first_page = len(text.strip()) > 10

                result.is_valid = (
                    page_count >= 1 and result.has_title_on_first_page
                )
                result.details["page_count"] = page_count
                result.details["lib_used"] = "PyPDF2"

        elif pdf_lib == "pdfplumber":
            import pdfplumber

            with pdfplumber.open(file_path) as pdf:
                page_count = len(pdf.pages)
                result.page_count = page_count

                # Check first page for content
                if page_count > 0:
                    first_page = pdf.pages[0]
                    text = first_page.extract_text() or ""
                    result.has_title_on_first_page = len(text.strip()) > 10

                result.is_valid = (
                    page_count >= 1 and result.has_title_on_first_page
                )
                result.details["page_count"] = page_count
                result.details["lib_used"] = "pdfplumber"

        logger.info(
            f"PDF layout check: {result.page_count} pages, "
            f"title_on_first_page={result.has_title_on_first_page}"
        )

    except Exception as e:
        logger.error(f"Error reading PDF: {e}")
        result.details["error"] = str(e)

    return result


def validate_layout_suite(
    html_path: Optional[str] = None,
    pptx_path: Optional[str] = None,
    pdf_path: Optional[str] = None,
) -> Dict[str, Any]:
    """
    Run all layout checks and return summary.

    Args:
        html_path: Path to HTML file (optional)
        pptx_path: Path to PPTX file (optional)
        pdf_path: Path to PDF file (optional)

    Returns:
        Dictionary with all layout check results
    """
    results = {}

    if html_path:
        results["html"] = check_html_layout(file_path=html_path)

    if pptx_path:
        results["pptx"] = check_pptx_layout(pptx_path)

    if pdf_path:
        results["pdf"] = check_pdf_layout(pdf_path)

    return results

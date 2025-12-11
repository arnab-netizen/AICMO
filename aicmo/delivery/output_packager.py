"""
Output Packager — Project Delivery Package Builder.

Phase 4: Bundle project outputs into deliverable packages.

This module packages completed project work into convenient formats:
- Generate strategy PDF from project strategy
- Generate creative/execution deck (PPTX) from creative outputs
- Generate HTML summary for web delivery
- Collate all deliverables with metadata
- Return structured ProjectPackageResult with paths and errors

Safety behavior:
- Each generator is wrapped in try/except; failures are recorded but don't stop processing
- Missing projects are handled gracefully
- If any file generation succeeds, result is marked as success=True
- Errors are logged but not raised from main function
- All operations are idempotent (safe to run multiple times)
"""

from dataclasses import dataclass, field
from typing import Optional, List, Dict, Any
import logging
import os
from datetime import datetime

from aicmo.memory.engine import log_event


logger = logging.getLogger(__name__)


# ═══════════════════════════════════════════════════════════════════════
# PACKAGE RESULT DATACLASS
# ═══════════════════════════════════════════════════════════════════════


@dataclass
class ProjectPackageResult:
    """
    Result of packaging project outputs into deliverable formats.
    
    Contains paths to generated files and metadata about the package.
    """
    
    project_id: str
    
    # Generated file paths
    pdf_path: Optional[str] = None
    pptx_path: Optional[str] = None
    html_summary_path: Optional[str] = None
    
    # Metadata
    metadata: Dict[str, Any] = field(default_factory=dict)
    
    # Status
    success: bool = False
    errors: List[str] = field(default_factory=list)
    
    created_at: Optional[str] = None
    
    def add_error(self, error_message: str) -> None:
        """Record an error in the package result."""
        if error_message not in self.errors:
            self.errors.append(error_message)
    
    def file_count(self) -> int:
        """Return count of generated files."""
        count = 0
        if self.pdf_path:
            count += 1
        if self.pptx_path:
            count += 1
        if self.html_summary_path:
            count += 1
        return count


# ═══════════════════════════════════════════════════════════════════════
# MAIN PACKAGER FUNCTION
# ═══════════════════════════════════════════════════════════════════════


def build_project_package(project_id: str) -> ProjectPackageResult:
    """
    Build complete deliverable package for a project.
    
    Generates PDF, PPTX, and HTML outputs from project data and stores
    references to generated files in the result. Each generation is
    wrapped in try/except to ensure partial success is useful.
    
    Safety behavior:
    - If project not found, returns result with success=False and error recorded
    - If any generator fails, continues with remaining generators
    - If any file is generated, success=True (partial success is ok)
    - Never raises exceptions from main function
    
    Args:
        project_id: Project UUID or ID string
        
    Returns:
        ProjectPackageResult with file paths and status
        
    Example:
        result = build_project_package("proj-123")
        if result.success:
            print(f"PDF: {result.pdf_path}")
            print(f"PPTX: {result.pptx_path}")
            print(f"HTML: {result.html_summary_path}")
        else:
            print(f"Errors: {result.errors}")
    """
    
    # Initialize result
    result = ProjectPackageResult(
        project_id=str(project_id),
        created_at=datetime.now().isoformat()
    )
    
    # Log package start
    log_event(
        "packaging.started",
        details={"project_id": str(project_id)},
        tags=["packaging", "phase4"]
    )
    
    # ─────────────────────────────────────────────────────────────
    # PHASE 1: Fetch project data
    # ─────────────────────────────────────────────────────────────
    project_data = None
    
    try:
        project_data = fetch_project_data(project_id)
        
        if project_data is None:
            error_msg = f"Project {project_id} not found"
            logger.warning(error_msg)
            log_event(
                "packaging.project_not_found",
                details={"project_id": str(project_id)},
                tags=["packaging", "error"]
            )
            result.add_error(error_msg)
            return result
        
        log_event(
            "packaging.project_fetched",
            details={"project_id": str(project_id)},
            tags=["packaging", "phase4"]
        )
    except Exception as e:
        error_msg = f"fetch_project_data failed: {str(e)}"
        logger.error(error_msg)
        log_event(
            "packaging.error",
            details={
                "stage": "fetch_project",
                "project_id": str(project_id),
                "error": str(e),
            },
            tags=["packaging", "error"]
        )
        result.add_error(error_msg)
        return result
    
    # ─────────────────────────────────────────────────────────────
    # PHASE 2: Generate strategy PDF
    # ─────────────────────────────────────────────────────────────
    try:
        pdf_path = generate_strategy_pdf(project_data)
        
        if pdf_path:
            result.pdf_path = pdf_path
            result.metadata["pdf_generated"] = True
            
            log_event(
                "packaging.pdf_generated",
                details={
                    "project_id": str(project_id),
                    "path": pdf_path,
                },
                tags=["packaging", "pdf"]
            )
        else:
            error_msg = "PDF generation returned None"
            logger.warning(error_msg)
            result.add_error(error_msg)
            
    except Exception as e:
        error_msg = f"generate_strategy_pdf failed: {str(e)}"
        logger.error(error_msg)
        log_event(
            "packaging.error",
            details={
                "stage": "pdf_generation",
                "project_id": str(project_id),
                "error": str(e),
            },
            tags=["packaging", "error"]
        )
        result.add_error(error_msg)
    
    # ─────────────────────────────────────────────────────────────
    # PHASE 3: Generate creative/execution deck (PPTX)
    # ─────────────────────────────────────────────────────────────
    try:
        pptx_path = generate_full_deck_pptx(project_data)
        
        if pptx_path:
            result.pptx_path = pptx_path
            result.metadata["pptx_generated"] = True
            
            log_event(
                "packaging.pptx_generated",
                details={
                    "project_id": str(project_id),
                    "path": pptx_path,
                },
                tags=["packaging", "pptx"]
            )
        else:
            error_msg = "PPTX generation returned None"
            logger.warning(error_msg)
            result.add_error(error_msg)
            
    except Exception as e:
        error_msg = f"generate_full_deck_pptx failed: {str(e)}"
        logger.error(error_msg)
        log_event(
            "packaging.error",
            details={
                "stage": "pptx_generation",
                "project_id": str(project_id),
                "error": str(e),
            },
            tags=["packaging", "error"]
        )
        result.add_error(error_msg)
    
    # ─────────────────────────────────────────────────────────────
    # PHASE 4: Generate HTML summary
    # ─────────────────────────────────────────────────────────────
    try:
        html_path = generate_html_summary(project_data)
        
        if html_path:
            result.html_summary_path = html_path
            result.metadata["html_generated"] = True
            
            log_event(
                "packaging.html_generated",
                details={
                    "project_id": str(project_id),
                    "path": html_path,
                },
                tags=["packaging", "html"]
            )
        else:
            error_msg = "HTML generation returned None"
            logger.warning(error_msg)
            result.add_error(error_msg)
            
    except Exception as e:
        error_msg = f"generate_html_summary failed: {str(e)}"
        logger.error(error_msg)
        log_event(
            "packaging.error",
            details={
                "stage": "html_generation",
                "project_id": str(project_id),
                "error": str(e),
            },
            tags=["packaging", "error"]
        )
        result.add_error(error_msg)
    
    # ─────────────────────────────────────────────────────────────
    # PHASE 5: Populate metadata and determine success
    # ─────────────────────────────────────────────────────────────
    
    # Success = any file generated
    result.success = result.file_count() > 0
    
    # Add summary metadata
    result.metadata["total_files"] = result.file_count()
    result.metadata["error_count"] = len(result.errors)
    result.metadata["files_generated"] = []
    
    if result.pdf_path:
        result.metadata["files_generated"].append("pdf")
    if result.pptx_path:
        result.metadata["files_generated"].append("pptx")
    if result.html_summary_path:
        result.metadata["files_generated"].append("html")
    
    # ─────────────────────────────────────────────────────────────
    # Build and return package result
    # ─────────────────────────────────────────────────────────────
    log_event(
        "packaging.completed",
        details={
            "project_id": str(project_id),
            "success": result.success,
            "files_generated": result.metadata.get("files_generated", []),
            "error_count": len(result.errors),
        },
        tags=["packaging", "phase4", "complete"]
    )
    
    return result


# ═══════════════════════════════════════════════════════════════════════
# HELPER FUNCTIONS (Generator Adapters)
# ═══════════════════════════════════════════════════════════════════════


def fetch_project_data(project_id: str) -> Optional[Dict[str, Any]]:
    """
    Fetch complete project data including strategy, creative, and execution.
    
    Args:
        project_id: Project UUID or ID
        
    Returns:
        Dictionary with project data, or None if not found
        
    Raises:
        Exception: On database or fetch errors
    """
    # TODO: Implement fetch from CampaignDB or Project store
    # Should retrieve:
    # - Project metadata
    # - Strategy document
    # - Creative assets
    # - Content calendar
    # - Execution results
    logger.debug(f"fetch_project_data: {project_id}")
    return None


def generate_strategy_pdf(project_data: Dict[str, Any]) -> Optional[str]:
    """
    Generate strategy PDF from project strategy document.
    
    Creates a formatted PDF of the strategy recommendations, market analysis,
    and implementation roadmap.
    
    Args:
        project_data: Project dictionary from fetch_project_data
        
    Returns:
        Path to generated PDF file, or None if generation failed
        
    Raises:
        Exception: On PDF generation errors
    """
    # TODO: Implement PDF generation
    # Should call underlying report generator or PDF library
    # Expected to exist as generate_strategy_pdf in reporting module
    logger.debug("generate_strategy_pdf")
    return None


def generate_full_deck_pptx(project_data: Dict[str, Any]) -> Optional[str]:
    """
    Generate creative/execution deck as PowerPoint PPTX.
    
    Creates a formatted presentation of creative assets, content calendar,
    and execution plan.
    
    Args:
        project_data: Project dictionary from fetch_project_data
        
    Returns:
        Path to generated PPTX file, or None if generation failed
        
    Raises:
        Exception: On PPTX generation errors
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from datetime import datetime
        import tempfile
        import os
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Title Slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        left = Inches(0.5)
        top = Inches(2.5)
        width = Inches(9)
        height = Inches(2)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = project_data.get("project_name", "Creative Execution Deck")
        title_frame.paragraphs[0].font.size = Pt(54)
        title_frame.paragraphs[0].font.bold = True
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = f"Generated: {datetime.now().strftime('%B %d, %Y')}"
        subtitle_frame.paragraphs[0].font.size = Pt(18)
        
        # Strategy Slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = "Strategy Overview"
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        strategy_text = project_data.get("strategy", "")[:500] or "Strategy details available"
        content_frame.text = strategy_text
        content_frame.paragraphs[0].font.size = Pt(14)
        content_frame.paragraphs[0].level = 0
        
        # Creatives Slide (per platform)
        platforms = project_data.get("platforms", [])
        for platform in platforms[:3]:  # Limit to 3 platform slides
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
            title_frame = title_box.text_frame
            title_frame.text = f"{platform.title()} Content"
            title_frame.paragraphs[0].font.size = Pt(40)
            title_frame.paragraphs[0].font.bold = True
            
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            platform_content = f"• Primary platform: {platform}\n• Content strategy tailored for {platform}\n• Expected engagement metrics"
            content_frame.text = platform_content
            content_frame.paragraphs[0].font.size = Pt(14)
        
        # Calendar Slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = "Content Calendar"
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        calendar_data = project_data.get("calendar", [])
        calendar_text = f"• Total posts planned: {len(calendar_data)}\n• Posting frequency: Daily\n• Best performing times: TBD"
        content_frame.text = calendar_text
        content_frame.paragraphs[0].font.size = Pt(14)
        
        # Save to temp file
        temp_dir = tempfile.gettempdir()
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        pptx_filename = f"creative_deck_{timestamp}.pptx"
        pptx_path = os.path.join(temp_dir, pptx_filename)
        
        prs.save(pptx_path)
        logger.info(f"Generated PPTX: {pptx_path}")
        return pptx_path
        
    except ImportError:
        logger.warning("python-pptx not installed, cannot generate PPTX")
        return None
    except Exception as e:
        logger.error(f"PPTX generation error: {e}")
        return None


def generate_html_summary(project_data: Dict[str, Any]) -> Optional[str]:
    """
    Generate HTML summary for web delivery.
    
    Creates a web-friendly HTML document with project overview,
    key insights, and deliverables.
    
    Args:
        project_data: Project dictionary from fetch_project_data
        
    Returns:
        Path to generated HTML file, or None if generation failed
        
    Raises:
        Exception: On HTML generation errors
    """
    try:
        from jinja2 import Template
        from datetime import datetime
        import tempfile
        import os
        import html as html_module
        
        # HTML template
        html_template = """
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{{ project_name }}</title>
    <style>
        * {
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }
        body {
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, Oxygen, Ubuntu, Cantarell, sans-serif;
            line-height: 1.6;
            color: #333;
            background: #f5f5f5;
        }
        .container {
            max-width: 900px;
            margin: 0 auto;
            padding: 40px 20px;
        }
        header {
            background: #2c3e50;
            color: white;
            padding: 40px 20px;
            margin: -40px -20px 40px -20px;
            border-radius: 8px 8px 0 0;
        }
        h1 {
            font-size: 2.5em;
            margin-bottom: 10px;
        }
        .meta {
            font-size: 0.9em;
            opacity: 0.9;
        }
        h2 {
            font-size: 1.8em;
            margin-top: 30px;
            margin-bottom: 15px;
            border-bottom: 3px solid #3498db;
            padding-bottom: 10px;
        }
        section {
            background: white;
            padding: 20px;
            margin-bottom: 20px;
            border-radius: 8px;
            box-shadow: 0 2px 4px rgba(0,0,0,0.1);
        }
        .platforms {
            display: flex;
            gap: 10px;
            margin-top: 15px;
            flex-wrap: wrap;
        }
        .platform-tag {
            background: #3498db;
            color: white;
            padding: 8px 16px;
            border-radius: 20px;
            font-size: 0.9em;
            font-weight: 600;
        }
        .calendar-preview {
            margin-top: 15px;
        }
        .calendar-item {
            background: #ecf0f1;
            padding: 10px;
            margin-bottom: 8px;
            border-left: 4px solid #3498db;
            border-radius: 4px;
        }
        footer {
            text-align: center;
            color: #7f8c8d;
            font-size: 0.9em;
            margin-top: 40px;
            padding-top: 20px;
            border-top: 1px solid #ecf0f1;
        }
    </style>
</head>
<body>
    <header>
        <h1>{{ project_name }}</h1>
        <div class="meta">Generated {{ generated_at }}</div>
    </header>
    
    <div class="container">
        <section>
            <h2>Project Overview</h2>
            <p>{{ overview }}</p>
        </section>
        
        <section>
            <h2>Strategy</h2>
            <p>{{ strategy }}</p>
            {% if platforms %}
            <div class="platforms">
                {% for platform in platforms %}
                <span class="platform-tag">{{ platform }}</span>
                {% endfor %}
            </div>
            {% endif %}
        </section>
        
        {% if calendar %}
        <section>
            <h2>Content Calendar</h2>
            <p>Total posts planned: {{ calendar|length }}</p>
            <div class="calendar-preview">
                {% for item in calendar[:7] %}
                <div class="calendar-item">
                    <strong>{{ item.date or 'TBD' }}</strong> - {{ item.platform or 'Multi' }}: {{ item.hook[:60] or 'Content post' }}...
                </div>
                {% endfor %}
                {% if calendar|length > 7 %}
                <p style="margin-top: 10px; color: #7f8c8d;"><em>+ {{ calendar|length - 7 }} more posts</em></p>
                {% endif %}
            </div>
        </section>
        {% endif %}
        
        <section>
            <h2>Deliverables</h2>
            <ul>
                <li>Content strategy document</li>
                <li>{{ calendar|length }} social media posts (planned)</li>
                <li>Multi-platform content calendar</li>
                <li>Monthly performance metrics framework</li>
            </ul>
        </section>
        
        <footer>
            <p>This report was automatically generated. All content is confidential.</p>
        </footer>
    </div>
</body>
</html>
        """
        
        # Prepare template data
        template_data = {
            "project_name": html_module.escape(project_data.get("project_name", "Marketing Strategy")),
            "generated_at": datetime.now().strftime("%B %d, %Y at %H:%M"),
            "overview": html_module.escape(project_data.get("overview", "")[:500] or "Project overview"),
            "strategy": html_module.escape(project_data.get("strategy", "")[:1000] or "Strategy details"),
            "platforms": project_data.get("platforms", []),
            "calendar": project_data.get("calendar", []),
        }
        
        # Render template
        template = Template(html_template)
        html_content = template.render(**template_data)
        
        # Save to temp file
        temp_dir = tempfile.gettempdir()
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        html_filename = f"summary_{timestamp}.html"
        html_path = os.path.join(temp_dir, html_filename)
        
        with open(html_path, "w", encoding="utf-8") as f:
            f.write(html_content)
        
        logger.info(f"Generated HTML summary: {html_path}")
        return html_path
        
    except Exception as e:
        logger.error(f"HTML generation error: {e}")
        return None
